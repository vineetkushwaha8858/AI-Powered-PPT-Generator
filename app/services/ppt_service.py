# from pptx import Presentation
# from pptx.util import Pt, Inches
# from pptx.enum.shapes import PP_PLACEHOLDER
# from .ai_service import FastUniversalAISystem
# from .content_service import FastContentGenerator
# from app.core.utils import clean_topic

# class FastPresentationGenerator:
#     def __init__(self):
#         self.ai_system = FastUniversalAISystem()
#         self.content_generator = FastContentGenerator(self.ai_system)

#     def generate_presentation_fast(self, user_topic, template_path, num_slides):
#         professional_topic = clean_topic(user_topic)
#         slides_content = self.content_generator.generate_fast_content(
#             professional_topic, max(num_slides-2, 0)
#         )
#         print("DEBUG SLIDES CONTENT:", slides_content)
#         final_slides_content = [{"type": "title_only", "title": professional_topic}]
#         for item in slides_content:
#             if "type" not in item:
#                 item["type"] = "bullet"
#             # SKIP BLANK bullet/content slides
#             if (
#                 item.get("type") in ["bullet", "text_content"]
#                 and not item.get("bullets", [])
#                 and not item.get("details", "").strip()
#             ):
#                 print("Skipping empty slide:", item)
#                 continue
#             final_slides_content.append(item)
#         # Always add Thank You slide last (no duplicate)
#         if not final_slides_content or final_slides_content[-1].get("type", "") != "thanks":
#             final_slides_content.append({"type": "thanks"})
#         print("DEBUG FINAL SLIDES:", final_slides_content)
#         return final_slides_content, professional_topic

# def remove_all_slides(prs):
#     for i in range(len(prs.slides) - 1, -1, -1):
#         rId = prs.slides._sldIdLst[i].rId
#         prs.part.drop_rel(rId)
#         del prs.slides._sldIdLst[i]

# def find_layout(template_info, slide_type):
#     priorities = {
#         "title_only": ["Title Slide", "Title Only"],
#         "bullet": ["Title and Content", "Title Only"],
#         "text_content": ["Title and Content", "Title Only"],
#         "thanks": ["Title Slide", "Title Only"]
#     }
#     for pref in priorities.get(slide_type, []):
#         for l in template_info["layouts"]:
#             if pref.lower() in l["layout_name"].lower():
#                 return l["layout_index"]
#     return {"title_only": 0, "bullet": 1, "text_content": 1, "thanks": 0}.get(slide_type, 1)

# def apply_formatting(slide, slide_data):
#     t = slide_data["type"].lower()
#     try:
#         if t == "title_only":
#             for shape in slide.shapes:
#                 ph_type = None
#                 if hasattr(shape, "placeholder_format"):
#                     try:
#                         ph_type = shape.placeholder_format.type
#                     except:
#                         pass
#                 if shape.has_text_frame and (ph_type == PP_PLACEHOLDER.TITLE or ph_type is None):
#                     shape.text_frame.clear()
#                     run = shape.text_frame.paragraphs[0].add_run()
#                     run.text = slide_data["title"]
#                     run.font.size = Pt(36)
#                     run.font.bold = True
#                     break
#         elif t == "bullet" or t == "text_content":
#             ph = 0
#             for shape in slide.shapes:
#                 if shape.has_text_frame and hasattr(shape, "placeholder_format"):
#                     if ph == 0:
#                         shape.text_frame.clear()
#                         run = shape.text_frame.paragraphs[0].add_run()
#                         run.text = slide_data.get("title", "")
#                         run.font.size = Pt(28)
#                         run.font.bold = True
#                     elif ph == 1:
#                         shape.text_frame.clear()
#                         bullets = slide_data.get("bullets", [])
#                         details = slide_data.get("details", "")
#                         if bullets:
#                             for i, bullet in enumerate(bullets):
#                                 clean_bullet = bullet.strip()
#                                 if clean_bullet:
#                                     if i == 0:
#                                         p = shape.text_frame.paragraphs[0]
#                                     else:
#                                         p = shape.text_frame.add_paragraph()
#                                     p.text = clean_bullet
#                                     p.level = 0
#                                     p.font.size = Pt(22)
#                         elif details:
#                             p = shape.text_frame.paragraphs[0]
#                             p.text = details
#                             p.level = 0
#                             p.font.size = Pt(22)
#                     ph += 1
#         elif t == "thanks":
#             text_added = False
#             for shape in slide.shapes:
#                 if shape.has_text_frame:
#                     shape.text_frame.clear()
#                     run = shape.text_frame.paragraphs[0].add_run()
#                     run.text = "Thank You\nQuestions & Discussion"
#                     run.font.size = Pt(28)
#                     run.font.bold = True
#                     text_added = True
#                     break
#             if not text_added:
#                 left = Inches(1.5)
#                 top = Inches(2.1)
#                 width = Inches(7.5)
#                 height = Inches(2.2)
#                 textbox = slide.shapes.add_textbox(left, top, width, height)
#                 tf = textbox.text_frame
#                 tf.clear()
#                 run = tf.paragraphs[0].add_run()
#                 run.text = "Thank You\nQuestions & Discussion"
#                 run.font.size = Pt(32)
#                 run.font.bold = True
#     except Exception as e:
#         print("Formatting Error:", e)
#         pass

# def generate_powerpoint(template_path, slides_content, template_info, topic):
#     prs = Presentation(template_path)
#     remove_all_slides(prs)
#     for slide_data in slides_content:
#         slide_type = slide_data.get("type", "bullet")
#         layout_idx = find_layout(template_info, slide_type)
#         slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
#         if slide_type == "title_only":
#             slide_data["title"] = topic
#         apply_formatting(slide, slide_data)
#     return prs



"""
Complete PPT Service with Spell Checking
File: app/services/ppt_service.py
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import PP_PLACEHOLDER
from io import BytesIO
from .ai_service import FastUniversalAISystem
from .content_service import FastContentGenerator
from .spell_checker import spell_checker  # ✅ Import spell checker
from app.core.utils import clean_topic


class FastPresentationGenerator:
    """Generate presentations with spell checking"""
    
    def __init__(self):
        self.ai_system = FastUniversalAISystem()
        self.content_generator = FastContentGenerator(self.ai_system)

    def generate_presentation_fast(self, user_topic, template_path, num_slides):
        """Generate with comprehensive spell checking"""
        
        # Clean and fix spelling in topic
        cleaned_topic = clean_topic(user_topic)
        professional_topic = spell_checker.fix_text(cleaned_topic)
        
        print(f"\n{'='*60}")
        print(f"📊 Generating Presentation")
        if cleaned_topic != professional_topic:
            print(f"   ✏️  Topic corrected:")
            print(f"      Before: {cleaned_topic}")
            print(f"      After:  {professional_topic}")
        else:
            print(f"   Topic: {professional_topic}")
        print(f"   Slides: {num_slides}")
        print(f"{'='*60}\n")
        
        # Generate content
        content_slides = max(num_slides - 2, 1)
        slides_content = self.content_generator.generate_fast_content(
            professional_topic, content_slides
        )
        
        print(f"✓ Generated {len(slides_content)} slides\n")
        
        # Build final slides with spell checking
        final_slides_content = []
        
        # 1. Title slide
        final_slides_content.append({
            "type": "title_only",
            "title": professional_topic
        })
        print(f"✓ Title slide")
        
        # 2. Content slides with spell checking
        valid_count = 0
        for idx, item in enumerate(slides_content):
            # Skip duplicate titles
            if item.get("type") == "title_only":
                print(f"⚠️  Skipped duplicate title #{idx+1}")
                continue
            
            if "type" not in item:
                item["type"] = "bullet"
            
            slide_type = item.get("type", "")
            has_bullets = bool(item.get("bullets", []))
            has_details = bool(item.get("details", "").strip())
            has_title = bool(item.get("title", "").strip())
            
            # Skip empty
            if slide_type in ["bullet", "text_content"]:
                if not has_bullets and not has_details:
                    print(f"⚠️  Skipped empty #{idx+1}")
                    continue
                
                if not has_title:
                    item["title"] = f"Point {valid_count + 1}"
                
                # ✅ APPLY SPELL CHECKING
                print(f"  📝 Checking slide {valid_count + 2}...")
                
                # Fix title
                original_title = item["title"]
                item["title"] = spell_checker.fix_text(original_title)
                if item["title"] != original_title:
                    print(f"    Title: '{original_title}' → '{item['title']}'")
                
                # Fix bullets
                if has_bullets:
                    corrected_bullets = []
                    for bullet in item["bullets"]:
                        if bullet:
                            original_bullet = str(bullet)
                            corrected_bullet = spell_checker.fix_text(original_bullet)
                            corrected_bullets.append(corrected_bullet)
                            
                            if corrected_bullet != original_bullet:
                                print(f"    Bullet: '{original_bullet}' → '{corrected_bullet}'")
                    
                    item["bullets"] = corrected_bullets
                
                # Fix details
                if has_details:
                    original_details = item["details"]
                    item["details"] = spell_checker.fix_text(original_details)
                    if item["details"] != original_details:
                        print(f"    Details corrected")
            
            final_slides_content.append(item)
            valid_count += 1
            print(f"✓ Slide {valid_count + 1} checked and added")
        
        # 3. Thank you slide
        if not final_slides_content or final_slides_content[-1].get("type") != "thanks":
            final_slides_content.append({"type": "thanks", "title": "Thank You"})
            print(f"✓ Thank you slide")
        
        print(f"\n📊 Total: {len(final_slides_content)} slides\n")
        
        return final_slides_content, professional_topic


def remove_all_slides(prs):
    """Remove existing slides"""
    count = len(prs.slides)
    for i in range(count - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    print(f"  ✓ Removed {count} template slides")


def find_layout(template_info, slide_type):
    """Find layout"""
    priorities = {
        "title_only": ["Title Slide", "Title Only"],
        "bullet": ["Title and Content", "Content"],
        "text_content": ["Title and Content"],
        "thanks": ["Title Slide", "Title Only"]
    }
    
    for pref in priorities.get(slide_type, []):
        for layout in template_info.get("layouts", []):
            if pref.lower() in layout["layout_name"].lower():
                return layout["layout_index"]
    
    return {"title_only": 0, "bullet": 1, "text_content": 1, "thanks": 0}.get(slide_type, 1)


def apply_formatting(slide, slide_data):
    """Apply formatting"""
    
    slide_type = slide_data.get("type", "").lower()
    
    try:
        print(f"    🔹 {slide_type}")
        
        if slide_type == "title_only":
            title_text = slide_data.get("title", "Untitled")
            title_shape = None
            to_remove = []
            
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    try:
                        ph_type = shape.placeholder_format.type
                        if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                            title_shape = shape
                        else:
                            to_remove.append(shape)
                    except:
                        pass
            
            if not title_shape and hasattr(slide.shapes, 'title'):
                title_shape = slide.shapes.title
            
            if title_shape and title_shape.has_text_frame:
                title_shape.text_frame.clear()
                p = title_shape.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = title_text
                run.font.size = Pt(44)
                run.font.bold = True
                print(f"      ✓ Title: {title_text[:40]}")
            
            # Remove other placeholders
            removed = 0
            for shape in to_remove:
                try:
                    sp = shape.element
                    sp.getparent().remove(sp)
                    removed += 1
                except:
                    try:
                        shape.left = Inches(-100)
                        shape.top = Inches(-100)
                        if shape.has_text_frame:
                            shape.text_frame.clear()
                        removed += 1
                    except:
                        pass
            
            if removed > 0:
                print(f"      ✓ Removed {removed} placeholder(s)")
        
        elif slide_type in ["bullet", "text_content"]:
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text_frame.clear()
                p = title_shape.text_frame.paragraphs[0]
                p.text = slide_data.get("title", "")
                p.font.bold = True
                p.font.size = Pt(28)
                print(f"      ✓ Title: {slide_data.get('title', '')[:40]}")
            
            content_shape = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape != slide.shapes.title:
                    if shape.has_text_frame:
                        try:
                            ph_type = shape.placeholder_format.type
                            if ph_type in [PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY]:
                                content_shape = shape
                                break
                        except:
                            content_shape = shape
                            break
            
            if content_shape and content_shape.has_text_frame:
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                bullets = slide_data.get("bullets", [])
                details = slide_data.get("details", "")
                
                if bullets:
                    count = 0
                    for bullet in bullets:
                        text = str(bullet).strip()
                        if not text:
                            continue
                        
                        p = text_frame.paragraphs[0] if count == 0 else text_frame.add_paragraph()
                        p.text = text
                        p.level = 0
                        p.font.size = Pt(18)
                        count += 1
                    
                    print(f"      ✓ {count} bullets")
                
                elif details:
                    p = text_frame.paragraphs[0]
                    p.text = details
                    p.font.size = Pt(18)
                    print(f"      ✓ Content")
        
        elif slide_type == "thanks":
            title_shape = None
            
            for shape in list(slide.shapes):
                if hasattr(shape, 'placeholder_format'):
                    try:
                        ph_type = shape.placeholder_format.type
                        if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                            title_shape = shape
                        else:
                            try:
                                sp = shape.element
                                sp.getparent().remove(sp)
                            except:
                                if shape.has_text_frame:
                                    shape.text_frame.clear()
                                try:
                                    shape.left = Inches(-100)
                                    shape.top = Inches(-100)
                                except:
                                    pass
                    except:
                        pass
            
            if not title_shape and hasattr(slide.shapes, 'title'):
                title_shape = slide.shapes.title
            
            if title_shape and title_shape.has_text_frame:
                title_shape.text_frame.clear()
                p = title_shape.text_frame.paragraphs[0]
                p.text = "Thank You"
                p.font.size = Pt(44)
                p.font.bold = True
                print(f"      ✓ Thank You title")
            
            try:
                textbox = slide.shapes.add_textbox(
                    left=Inches(1.5),
                    top=Inches(5.0),
                    width=Inches(7),
                    height=Inches(1.2)
                )
                
                tf = textbox.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = "Questions & Discussion"
                p.font.size = Pt(26)
                p.alignment = 1
                
                print(f"      ✓ Subtitle positioned below")
                
            except Exception as e:
                print(f"      ⚠️  Subtitle error: {e}")
    
    except Exception as e:
        print(f"      ❌ Error: {e}")


def generate_powerpoint(template_path, slides_content, template_info, topic):
    """Generate PowerPoint"""
    
    print(f"\n{'='*60}")
    print(f"🎨 Creating PowerPoint")
    print(f"{'='*60}")
    print(f"Template: {template_path}")
    print(f"Topic: {topic}")
    print(f"Slides: {len(slides_content)}")
    print(f"{'='*60}\n")
    
    prs = Presentation(template_path)
    
    print(f"📐 Layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"   {i}: {layout.name}")
    print()
    
    remove_all_slides(prs)
    print()
    
    for idx, slide_data in enumerate(slides_content):
        print(f"  Slide {idx + 1}/{len(slides_content)}:")
        
        slide_type = slide_data.get("type", "bullet")
        layout_idx = find_layout(template_info, slide_type)
        
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            
            if slide_type == "title_only":
                slide_data["title"] = topic
            
            apply_formatting(slide, slide_data)
            print(f"    ✅ Done\n")
            
        except Exception as e:
            print(f"    ❌ Failed: {e}\n")
            continue
    
    print(f"{'='*60}")
    print(f"✅ Created: {len(prs.slides)} slides")
    print(f"{'='*60}\n")
    
    return prs


def save_presentation(prs, filename=None):
    """Save presentation"""
    if filename:
        prs.save(filename)
        print(f"💾 Saved: {filename}")
        return filename
    else:
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        print(f"💾 Saved: {output.getbuffer().nbytes:,} bytes")
        return output
